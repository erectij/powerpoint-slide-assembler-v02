"""
PowerPoint Slide Assembler - Streamlit Application - ENHANCED VERSION
=====================================================================

This application combines:
1. Template duplication method for perfect formatting preservation
2. DeepSeek LLM integration for intelligent slide analysis
3. User-friendly Streamlit interface
4. Comprehensive logging and debugging

**ENHANCED v2.1**: 
- Fixed nested expander errors
- Added comprehensive timeout handling and debugging
- Extended timeouts for large content (3 minutes)
- Content size optimization to prevent timeouts
- Detailed progress tracking and error analysis
- Better performance monitoring and token usage tracking
- SECURITY: API key now uses Streamlit secrets
- PRIVACY: Added comprehensive privacy warning system

Requirements:
    pip install streamlit python-pptx openai httpx requests

Usage:
    streamlit run app.py

Author: AI Assistant
Date: 2025-01-06
Version: 2.1 - Enhanced debugging, timeout fixes, content optimization, security & privacy
"""

import streamlit as st
import os
import logging
import datetime
import shutil
import tempfile
import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from openai import OpenAI
import time
import traceback
from typing import List, Dict, Tuple, Optional

# ============================================================================
# CONFIGURATION AND SETUP
# ============================================================================

# DeepSeek API Configuration via OpenRouter
OPENROUTER_BASE_URL = "https://openrouter.ai/api/v1"
DEEPSEEK_MODEL = "deepseek/deepseek-r1-0528:free"

# Application Configuration
APP_VERSION = "2.1-ENHANCED-SECURE"
APP_NAME = "PowerPoint Slide Assembler"
DEBUG_MODE = True  # Set to True for detailed debugging output

# File size limits
MAX_FILE_SIZE_MB = 100
MAX_TOTAL_FILES = 10

# Output configuration
OUTPUT_DIR_NAME = "slide_assembler_output"

# ============================================================================
# LOGGING SETUP
# ============================================================================

def setup_comprehensive_logging(session_id: str) -> Tuple[Path, logging.Logger]:
    """Setup comprehensive logging system with both file and console output."""
    # Create output directory structure
    output_dir = Path(OUTPUT_DIR_NAME)
    output_dir.mkdir(exist_ok=True)
    
    logs_dir = output_dir / "logs"
    logs_dir.mkdir(exist_ok=True)
    
    # Create timestamped log file
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_file = logs_dir / f"slide_assembler_{session_id}_{timestamp}.log"
    
    # Configure logger
    logger = logging.getLogger(f"slide_assembler_{session_id}")
    logger.setLevel(logging.DEBUG if DEBUG_MODE else logging.INFO)
    
    # Clear any existing handlers
    for handler in logger.handlers[:]:
        logger.removeHandler(handler)
    
    # File handler with detailed formatting
    file_handler = logging.FileHandler(log_file, encoding='utf-8')
    file_handler.setLevel(logging.DEBUG)
    file_formatter = logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(funcName)s:%(lineno)d - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    file_handler.setFormatter(file_formatter)
    logger.addHandler(file_handler)
    
    # Console handler for Streamlit (less verbose)
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)
    console_formatter = logging.Formatter(
        '%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%H:%M:%S'
    )
    console_handler.setFormatter(console_formatter)
    logger.addHandler(console_handler)
    
    # Log session start
    logger.info("=" * 80)
    logger.info(f"{APP_NAME} v{APP_VERSION} - Session Started")
    logger.info(f"Session ID: {session_id}")
    logger.info(f"Timestamp: {datetime.datetime.now().isoformat()}")
    logger.info(f"Debug Mode: {'ENABLED' if DEBUG_MODE else 'DISABLED'}")
    logger.info(f"Output Directory: {output_dir.absolute()}")
    logger.info(f"Log File: {log_file.name}")
    logger.info("=" * 80)
    
    return log_file, logger

# ============================================================================
# PRIVACY AND SECURITY FUNCTIONS
# ============================================================================

def show_privacy_warning():
    """Show privacy warning modal when files are uploaded."""
    
    # Initialize privacy acknowledgment in session state
    if 'privacy_acknowledged' not in st.session_state:
        st.session_state.privacy_acknowledged = False
    
    # Show warning if not yet acknowledged
    if not st.session_state.privacy_acknowledged:
        
        # Create a prominent warning container
        with st.container():
            st.error("🚨 **PRIVACY & SECURITY WARNING**")
            
            warning_text = """
**This application sends your presentation content to external AI services (OpenRouter/DeepSeek) for analysis.**

### 🚫 **DO NOT UPLOAD:**
- Confidential company information
- Personal data or PII (names, emails, addresses, phone numbers)
- Proprietary business content or trade secrets
- Internal strategies, financial data, or competitive information
- Client information or customer data
- Any content marked as confidential, internal, or restricted

### ✅ **SAFE TO UPLOAD:**
- Public presentation materials
- Educational or training content
- Non-sensitive examples and templates
- Marketing materials already public
- General business presentations without sensitive details

### ⚖️ **Legal Notice:**
By uploading files, you acknowledge that:
- Content will be processed by external AI services
- You have the right to share this content externally
- You take full responsibility for the content you upload
- No confidential or sensitive information is included

**Your data may be processed outside your organization's network and jurisdiction.**
            """
            
            st.markdown(warning_text)
            
            st.warning("**🔒 This warning protects your organization's confidential information. Please read carefully before proceeding.**")
            
            # Create acknowledgment section
            col1, col2, col3 = st.columns([1, 2, 1])
            
            with col2:
                # Checkbox for acknowledgment
                acknowledge_checkbox = st.checkbox(
                    "✅ **I understand the privacy implications and confirm that I will NOT upload any confidential or sensitive information**",
                    key="privacy_acknowledge_checkbox"
                )
                
                # Confirmation button
                if acknowledge_checkbox:
                    if st.button("🔓 **I Acknowledge - Proceed with File Upload**", type="primary", use_container_width=True):
                        st.session_state.privacy_acknowledged = True
                        st.session_state.logger.info("User acknowledged privacy warning")
                        st.rerun()
                else:
                    st.button("⚠️ **Please check the acknowledgment above to continue**", disabled=True, use_container_width=True)
        
        return False  # Block further functionality
    
    return True  # Privacy acknowledged, proceed

# ============================================================================
# LLM INTEGRATION - DEEPSEEK VIA OPENROUTER
# ============================================================================

def initialize_llm_client(api_key: str, bypass_ssl: bool = False) -> OpenAI:
    """Initialize OpenRouter client for DeepSeek LLM access."""
    logger = st.session_state.logger
    
    try:
        logger.info("Initializing OpenRouter client for DeepSeek API...")
        logger.debug(f"Base URL: {OPENROUTER_BASE_URL}")
        logger.debug(f"Model: {DEEPSEEK_MODEL}")
        logger.debug(f"API Key prefix: {api_key[:10]}..." if api_key else "No API key provided")
        logger.debug(f"SSL Bypass: {bypass_ssl}")
        
        # Validate API key format
        if not api_key:
            raise ValueError("API key is required")
        
        if not api_key.startswith("sk-or-v1-"):
            raise ValueError("Invalid OpenRouter API key format. Should start with 'sk-or-v1-'")
        
        # Configure HTTP client for SSL issues with longer timeout for large content
        try:
            import httpx
        except ImportError:
            st.error("Missing httpx dependency. Please install with: pip install httpx")
            return None
        
        # Increased timeout for large content analysis (3 minutes)
        timeout_config = httpx.Timeout(
            connect=30.0,   # Connection timeout
            read=180.0,     # Read timeout (3 minutes for large analysis)
            write=30.0,     # Write timeout
            pool=30.0       # Pool timeout
        )
        
        if bypass_ssl:
            logger.warning("⚠️ SSL verification disabled for corporate network compatibility")
            http_client = httpx.Client(verify=False, timeout=timeout_config)
        else:
            http_client = httpx.Client(timeout=timeout_config)
        
        logger.debug(f"HTTP client configured with timeouts: connect=30s, read=180s, write=30s")
        
        # Create client with custom HTTP client
        client = OpenAI(
            base_url=OPENROUTER_BASE_URL,
            api_key=api_key,
            http_client=http_client
        )
        
        logger.info("OpenRouter client created, testing connection...")
        
        # Test connection with a simple query
        test_completion = client.chat.completions.create(
            extra_headers={
                "HTTP-Referer": "https://slide-assembler.internal",
                "X-Title": "PowerPoint Slide Assembler",
            },
            extra_body={},
            model=DEEPSEEK_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": "Hello! Please respond with 'API connection successful' to confirm the connection."
                }
            ],
            max_tokens=50,
            temperature=0.1
        )
        
        response_text = test_completion.choices[0].message.content
        logger.info(f"✅ LLM API test successful: {response_text}")
        
        # Log token usage if available
        if hasattr(test_completion, 'usage'):
            usage = test_completion.usage
            logger.debug(f"Test API call token usage: {usage.total_tokens} total "
                        f"({usage.prompt_tokens} prompt + {usage.completion_tokens} completion)")
        
        return client
        
    except Exception as e:
        error_type = type(e).__name__
        error_message = str(e)
        
        logger.error(f"❌ Failed to initialize LLM client: {error_type}: {error_message}")
        
        # Provide specific error messages for common issues
        if "CERTIFICATE_VERIFY_FAILED" in error_message or "certificate verify failed" in error_message:
            user_message = """🔒 **Corporate SSL Certificate Issue**

Your corporate network is intercepting HTTPS traffic with its own certificates.

**Solutions:**
1. **Try SSL Bypass** (enable checkbox and reconnect)
2. **Contact IT** to whitelist `openrouter.ai` 
3. **Use personal hotspot** to test outside corporate network

This is very common in enterprise environments."""
            
        elif "Connection error" in error_message or "ConnectTimeout" in error_message:
            user_message = "🌐 **Network Connection Error**\n\nPossible causes:\n• Check your internet connection\n• Corporate firewall may be blocking OpenRouter\n• OpenRouter service may be temporarily unavailable"
        elif "401" in error_message or "Unauthorized" in error_message:
            user_message = "🔑 **Authentication Error**\n\nPossible causes:\n• Invalid API key\n• API key may be expired\n• Check that the key starts with 'sk-or-v1-'"
        elif "404" in error_message or "Not Found" in error_message:
            user_message = f"🤖 **Model Not Found**\n\nThe model '{DEEPSEEK_MODEL}' may not be available.\n• Check if the model name is correct\n• Model may be temporarily unavailable"
        elif "rate" in error_message.lower() or "limit" in error_message.lower():
            user_message = "⏱️ **Rate Limit Exceeded**\n\nToo many requests to the API.\n• Wait a few minutes and try again\n• Consider upgrading your OpenRouter plan"
        else:
            user_message = f"❌ **API Connection Failed**\n\nError: {error_message}\n\nPlease check:\n• API key is correct\n• Internet connection is working\n• OpenRouter service status"
        
        st.error(user_message)
        
        # Debug information
        if DEBUG_MODE:
            import traceback
            logger.debug(f"Full traceback: {traceback.format_exc()}")
            
            st.markdown("**🐛 Debug Information:**")
            st.code(f"Error Type: {error_type}\nError Message: {error_message}\n\nFull Traceback:\n{traceback.format_exc()}")
        
        return None

def analyze_slide_content_with_llm(client: OpenAI, slide_content: str, key_message: str, 
                                   analysis_type: str = "comprehensive") -> Dict:
    """Analyze slide content using DeepSeek LLM for intelligent recommendations."""
    logger = st.session_state.logger
    
    try:
        logger.info(f"Starting LLM analysis - Type: {analysis_type}")
        logger.info(f"Content length: {len(slide_content):,} characters")
        logger.info(f"Key message length: {len(key_message)} characters")
        logger.debug(f"Key message: {key_message}")
        
        # Check content size and warn if very large
        if len(slide_content) > 50000:
            logger.warning(f"⚠️ Very large content detected: {len(slide_content):,} characters")
            logger.warning("This may take longer to process or hit API limits")
        elif len(slide_content) > 20000:
            logger.info(f"Large content detected: {len(slide_content):,} characters - expect longer processing time")
        
        # Content preview for debugging
        logger.debug(f"Content preview (first 500 chars): {slide_content[:500]}...")
        
        # Construct analysis prompt based on type
        if analysis_type == "comprehensive":
            system_prompt = """You are an expert presentation consultant. Analyze the provided slide content and provide comprehensive recommendations for creating an effective presentation."""
            
            user_prompt = f"""
I need to create a presentation with the key message: "{key_message}"

Here is the content from slides I'm considering:

{slide_content}

Please provide a comprehensive analysis including:

1. **Relevance Assessment**: How well does each piece of content support the key message?
2. **Content Quality**: Evaluate the clarity, impact, and professionalism of the content
3. **Logical Flow**: Suggest the best order for presenting this information
4. **Missing Elements**: What key points or supporting information might be missing?
5. **Slide Recommendations**: Which slides are essential, nice-to-have, or should be excluded?
6. **Narrative Structure**: How to weave this content into a compelling story

Please be specific and actionable in your recommendations.
"""
        
        elif analysis_type == "relevance":
            system_prompt = """You are a presentation expert. Focus on evaluating content relevance and providing clear recommendations."""
            
            user_prompt = f"""
Key message: "{key_message}"

Slide content:
{slide_content}

Please evaluate:
1. Which slides directly support the key message?
2. Which slides are tangentially related but could be useful?
3. Which slides don't align and should be excluded?
4. Suggest a priority ranking for the relevant slides.

Be concise and specific.
"""
        
        else:  # quick analysis
            system_prompt = """You are a presentation consultant. Provide quick, actionable insights."""
            
            user_prompt = f"""
Key message: "{key_message}"

Slide content:
{slide_content}

Quick assessment:
1. Top 3 most relevant slides for this message
2. Overall content quality (1-10)
3. One key recommendation for improvement

Keep it brief and actionable.
"""
        
        # Make API call to DeepSeek with detailed logging
        logger.info("🚀 Preparing API request to DeepSeek...")
        logger.debug(f"System prompt length: {len(system_prompt)} characters")
        logger.debug(f"User prompt length: {len(user_prompt)} characters")
        logger.debug(f"Total prompt size: {len(system_prompt) + len(user_prompt):,} characters")
        
        # Estimate tokens (rough estimate: 1 token ≈ 4 characters)
        estimated_tokens = (len(system_prompt) + len(user_prompt)) // 4
        logger.info(f"Estimated input tokens: ~{estimated_tokens:,}")
        
        if estimated_tokens > 100000:
            logger.warning("⚠️ Very large prompt detected - may hit model limits or take very long")
        elif estimated_tokens > 50000:
            logger.warning("⚠️ Large prompt - expect longer processing time")
        
        # Log request parameters
        logger.debug(f"Request parameters:")
        logger.debug(f"  - Model: {DEEPSEEK_MODEL}")
        logger.debug(f"  - Max tokens: 2000")
        logger.debug(f"  - Temperature: 0.1")
        logger.debug(f"  - Messages count: 2 (system + user)")
        
        logger.info("📡 Sending request to DeepSeek API...")
        start_time = time.time()
        
        try:
            # Add progress logging for long requests
            logger.info("⏳ API request in progress... (this may take 1-3 minutes for large content)")
            
            completion = client.chat.completions.create(
                extra_headers={
                    "HTTP-Referer": "https://slide-assembler.internal",
                    "X-Title": "PowerPoint Slide Assembler - Content Analysis",
                },
                model=DEEPSEEK_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_tokens=2000,
                temperature=0.1
            )
            
            end_time = time.time()
            response_time = end_time - start_time
            
            logger.info(f"✅ API request completed successfully!")
            logger.info(f"📊 Response time: {response_time:.2f} seconds")
            
        except Exception as api_error:
            end_time = time.time()
            response_time = end_time - start_time
            
            logger.error(f"❌ API request failed after {response_time:.2f} seconds")
            logger.error(f"Error type: {type(api_error).__name__}")
            logger.error(f"Error message: {str(api_error)}")
            
            # Specific timeout handling
            if "timeout" in str(api_error).lower() or "timed out" in str(api_error).lower():
                logger.error("🕐 Request timed out - content may be too large or service overloaded")
                logger.info("💡 Suggestions:")
                logger.info("  - Try 'quick' or 'relevance' analysis instead of 'comprehensive'")
                logger.info("  - Reduce content by selecting fewer slides")
                logger.info("  - Try again in a few minutes (service may be busy)")
                
            # Re-raise the exception to be handled by outer try-catch
            raise api_error
        
        # Extract and process response with detailed logging
        logger.info("📝 Processing API response...")
        
        try:
            analysis_text = completion.choices[0].message.content
            logger.info(f"✅ Response extracted successfully")
            logger.info(f"Response length: {len(analysis_text):,} characters")
            logger.debug(f"Response preview (first 200 chars): {analysis_text[:200]}...")
            
        except (IndexError, AttributeError) as response_error:
            logger.error(f"❌ Failed to extract response content: {response_error}")
            logger.debug(f"Raw completion object: {completion}")
            raise ValueError(f"Invalid API response structure: {response_error}")
        
        # Calculate token usage (if available) with detailed logging
        token_usage = getattr(completion, 'usage', None)
        if token_usage:
            logger.info(f"📊 Token usage details:")
            logger.info(f"  - Prompt tokens: {token_usage.prompt_tokens:,}")
            logger.info(f"  - Completion tokens: {token_usage.completion_tokens:,}")
            logger.info(f"  - Total tokens: {token_usage.total_tokens:,}")
            
            # Calculate cost estimate (rough estimate for OpenRouter)
            estimated_cost = (token_usage.prompt_tokens * 0.00001) + (token_usage.completion_tokens * 0.00001)
            logger.debug(f"  - Estimated cost: ~${estimated_cost:.4f}")
        else:
            logger.warning("⚠️ Token usage information not available in response")
        
        # Log response metadata
        if hasattr(completion, 'id'):
            logger.debug(f"Request ID: {completion.id}")
        if hasattr(completion, 'model'):
            logger.debug(f"Model used: {completion.model}")
        if hasattr(completion, 'created'):
            logger.debug(f"Created timestamp: {completion.created}")
        
        # Create structured response
        analysis_result = {
            "success": True,
            "analysis_type": analysis_type,
            "analysis_text": analysis_text,
            "response_time": response_time,
            "content_stats": {
                "input_characters": len(slide_content),
                "output_characters": len(analysis_text),
                "estimated_input_tokens": estimated_tokens,
            },
            "token_usage": {
                "prompt_tokens": token_usage.prompt_tokens if token_usage else None,
                "completion_tokens": token_usage.completion_tokens if token_usage else None,
                "total_tokens": token_usage.total_tokens if token_usage else None,
            } if token_usage else None,
            "timestamp": datetime.datetime.now().isoformat(),
            "model": DEEPSEEK_MODEL
        }
        
        logger.info(f"🎉 LLM analysis completed successfully!")
        logger.info(f"📈 Performance summary:")
        logger.info(f"  - Processing time: {response_time:.2f} seconds")
        logger.info(f"  - Input: {len(slide_content):,} chars → Output: {len(analysis_text):,} chars")
        if token_usage:
            logger.info(f"  - Tokens used: {token_usage.total_tokens:,}")
        
        return analysis_result
        
    except Exception as e:
        error_type = type(e).__name__
        error_message = str(e)
        
        logger.error(f"❌ LLM analysis failed with {error_type}: {error_message}")
        logger.debug(f"Full error traceback: {traceback.format_exc()}")
        
        # Detailed error analysis and suggestions
        if "timeout" in error_message.lower() or "timed out" in error_message.lower():
            logger.error("🕐 TIMEOUT ERROR - Request took too long to complete")
            logger.info("📊 Timeout analysis:")
            logger.info(f"  - Content size: {len(slide_content):,} characters")
            logger.info(f"  - Estimated tokens: ~{(len(slide_content) // 4):,}")
            logger.info(f"  - Analysis type: {analysis_type}")
            logger.info("💡 Timeout solutions:")
            logger.info("  1. Try 'quick' analysis instead of 'comprehensive'")
            logger.info("  2. Select fewer slides to reduce content size")
            logger.info("  3. Break content into smaller chunks")
            logger.info("  4. Try again later (API may be overloaded)")
            
        elif "rate" in error_message.lower() or "limit" in error_message.lower():
            logger.error("🚦 RATE LIMIT ERROR - Too many requests")
            logger.info("💡 Rate limit solutions:")
            logger.info("  1. Wait 1-2 minutes before trying again")
            logger.info("  2. Consider upgrading OpenRouter plan")
            logger.info("  3. Use 'quick' analysis to reduce token usage")
            
        elif "401" in error_message or "unauthorized" in error_message.lower():
            logger.error("🔑 AUTHENTICATION ERROR - API key issue")
            logger.info("💡 Auth error solutions:")
            logger.info("  1. Check API key is correct")
            logger.info("  2. Verify API key hasn't expired")
            logger.info("  3. Ensure sufficient credits on OpenRouter account")
            
        elif "400" in error_message or "bad request" in error_message.lower():
            logger.error("📝 REQUEST ERROR - Invalid request format")
            logger.info("💡 Request error analysis:")
            logger.info(f"  - Content may be too large: {len(slide_content):,} chars")
            logger.info(f"  - Try reducing content size or using 'quick' analysis")
            
        elif "502" in error_message or "503" in error_message or "server" in error_message.lower():
            logger.error("🖥️ SERVER ERROR - OpenRouter/DeepSeek service issue")
            logger.info("💡 Server error solutions:")
            logger.info("  1. Try again in 2-3 minutes")
            logger.info("  2. Check OpenRouter status page")
            logger.info("  3. Consider using different analysis type")
            
        else:
            logger.error(f"❓ UNKNOWN ERROR - {error_type}")
            logger.info("💡 General troubleshooting:")
            logger.info("  1. Check internet connection")
            logger.info("  2. Try again with smaller content")
            logger.info("  3. Verify API configuration")
        
        # Log detailed error context
        logger.debug("🐛 Error context:")
        logger.debug(f"  - Analysis type: {analysis_type}")
        logger.debug(f"  - Content length: {len(slide_content):,}")
        logger.debug(f"  - Key message length: {len(key_message)}")
        logger.debug(f"  - Model: {DEEPSEEK_MODEL}")
        
        return {
            "success": False,
            "error": error_message,
            "error_type": error_type,
            "analysis_type": analysis_type,
            "content_stats": {
                "input_characters": len(slide_content),
                "estimated_input_tokens": len(slide_content) // 4,
            },
            "timestamp": datetime.datetime.now().isoformat(),
            "troubleshooting": {
                "content_size": "large" if len(slide_content) > 20000 else "normal",
                "suggested_actions": [
                    "Try 'quick' analysis type",
                    "Select fewer slides",
                    "Wait and retry",
                    "Check API status"
                ]
            }
        }

# ============================================================================
# SLIDE CONTENT EXTRACTION
# ============================================================================

def extract_slide_text_content(presentation_path: str, slide_indices: List[int] = None) -> Dict:
    """Extract text content from PowerPoint slides for LLM analysis."""
    logger = st.session_state.logger
    
    try:
        logger.debug(f"Extracting content from: {os.path.basename(presentation_path)}")
        
        # Open presentation
        presentation = Presentation(presentation_path)
        total_slides = len(presentation.slides)
        
        # Determine which slides to process
        if slide_indices is None:
            slides_to_process = list(range(total_slides))
        else:
            slides_to_process = [i for i in slide_indices if 0 <= i < total_slides]
        
        logger.debug(f"Processing {len(slides_to_process)} slides out of {total_slides} total")
        
        # Extract content from each slide
        slide_contents = []
        extraction_stats = {
            "total_slides": total_slides,
            "processed_slides": len(slides_to_process),
            "text_shapes": 0,
            "total_characters": 0,
            "extraction_time": 0
        }
        
        start_time = time.time()
        
        for slide_idx in slides_to_process:
            try:
                slide = presentation.slides[slide_idx]
                slide_text_parts = []
                slide_shape_count = 0
                
                logger.debug(f"Processing slide {slide_idx + 1}: {len(slide.shapes)} shapes")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())
                            slide_shape_count += 1
                            extraction_stats["text_shapes"] += 1
                            
                    except Exception as shape_error:
                        logger.debug(f"Error extracting from shape in slide {slide_idx + 1}: {shape_error}")
                        continue
                
                # Extract slide notes if available
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_text_parts.append(f"[NOTES: {notes_text}]")
                except Exception as notes_error:
                    logger.debug(f"Error extracting notes from slide {slide_idx + 1}: {notes_error}")
                
                # Combine slide content
                slide_content = "\\n\\n".join(slide_text_parts)
                
                slide_info = {
                    "slide_number": slide_idx + 1,
                    "slide_index": slide_idx,
                    "content": slide_content,
                    "character_count": len(slide_content),
                    "shape_count": slide_shape_count,
                    "has_content": bool(slide_content.strip())
                }
                
                slide_contents.append(slide_info)
                extraction_stats["total_characters"] += len(slide_content)
                
                logger.debug(f"Slide {slide_idx + 1}: {len(slide_content)} chars, {slide_shape_count} text shapes")
                
            except Exception as slide_error:
                logger.warning(f"Error processing slide {slide_idx + 1}: {slide_error}")
                continue
        
        extraction_stats["extraction_time"] = time.time() - start_time
        
        # Create combined content for LLM analysis
        combined_content = ""
        for slide_info in slide_contents:
            if slide_info["has_content"]:
                combined_content += f"\\n\\n=== SLIDE {slide_info['slide_number']} ===\\n"
                combined_content += slide_info["content"]
        
        result = {
            "success": True,
            "file_name": os.path.basename(presentation_path),
            "slide_contents": slide_contents,
            "combined_content": combined_content.strip(),
            "extraction_stats": extraction_stats,
            "timestamp": datetime.datetime.now().isoformat()
        }
        
        logger.info(f"Content extraction completed: {extraction_stats['processed_slides']} slides, "
                   f"{extraction_stats['total_characters']} characters, "
                   f"{extraction_stats['text_shapes']} text shapes")
        
        return result
        
    except Exception as e:
        logger.error(f"Content extraction failed: {e}")
        return {
            "success": False,
            "error": str(e),
            "file_name": os.path.basename(presentation_path) if presentation_path else "Unknown",
            "timestamp": datetime.datetime.now().isoformat()
        }

# ============================================================================
# TEMPLATE DUPLICATION METHODS (FROM PROVEN APPROACH)
# ============================================================================

def parse_slide_ranges(range_string: str) -> List[int]:
    """Parse slide range string like "1-3,6,9-12" into list of slide numbers."""
    logger = st.session_state.logger
    slide_numbers = []
    
    try:
        logger.debug(f"Parsing slide range string: '{range_string}'")
        
        # Split by commas and process each part
        parts = [part.strip() for part in range_string.split(',') if part.strip()]
        
        for part in parts:
            if '-' in part:
                # Handle ranges like "1-3" or "9-12"
                try:
                    start, end = part.split('-', 1)
                    start_num = int(start.strip())
                    end_num = int(end.strip())
                    
                    if start_num > end_num:
                        logger.warning(f"Invalid range {part}: start > end")
                        continue
                    
                    range_list = list(range(start_num, end_num + 1))
                    slide_numbers.extend(range_list)
                    logger.debug(f"Parsed range '{part}' -> {range_list}")
                    
                except ValueError as ve:
                    logger.error(f"Invalid range format '{part}': {ve}")
                    continue
                    
            else:
                # Handle single numbers like "6"
                try:
                    num = int(part.strip())
                    if num > 0:  # Ensure positive slide numbers
                        slide_numbers.append(num)
                        logger.debug(f"Parsed single number: {num}")
                    else:
                        logger.warning(f"Invalid slide number {num}: must be positive")
                        
                except ValueError as ve:
                    logger.error(f"Invalid number '{part}': {ve}")
                    continue
        
        # Remove duplicates and sort
        slide_numbers = sorted(list(set(slide_numbers)))
        logger.info(f"Final parsed slide numbers: {slide_numbers}")
        
    except Exception as e:
        logger.error(f"Error parsing slide ranges '{range_string}': {e}")
        return []
    
    return slide_numbers

def create_trimmed_template(source_file_path: str, slides_to_keep: List[int], 
                          output_dir: Path, file_identifier: str) -> Dict:
    """Create a trimmed template using the proven template duplication method."""
    logger = st.session_state.logger
    
    # Initialize result structure
    result = {
        "success": False,
        "template_path": None,
        "source_file": os.path.basename(source_file_path),
        "slides_requested": slides_to_keep.copy(),
        "original_slide_count": 0,
        "final_slide_count": 0,
        "slides_deleted": 0,
        "processing_time": 0,
        "error": None,
        "timestamp": datetime.datetime.now().isoformat()
    }
    
    start_time = time.time()
    
    try:
        logger.info(f"Creating trimmed template from: {os.path.basename(source_file_path)}")
        logger.info(f"Slides to keep: {slides_to_keep}")
        
        # Create temp directory if it doesn't exist
        temp_dir = output_dir / "temp_templates"
        temp_dir.mkdir(exist_ok=True)
        
        # Generate unique temp file name
        timestamp = datetime.datetime.now().strftime("%H%M%S_%f")[:-3]
        temp_file = temp_dir / f"template_{file_identifier}_{timestamp}.pptx"
        
        # Step 1: Copy source file to temp location (preserves everything!)
        logger.debug(f"Copying source file to: {temp_file.name}")
        shutil.copy2(source_file_path, temp_file)
        
        # Step 2: Open the copied presentation for modification
        logger.debug("Opening copied presentation for modification")
        presentation = Presentation(str(temp_file))
        result["original_slide_count"] = len(presentation.slides)
        
        logger.info(f"Original presentation has {result['original_slide_count']} slides")
        
        # Step 3: Validate that all requested slides exist
        max_slides = result["original_slide_count"]
        invalid_slides = [s for s in slides_to_keep if s < 1 or s > max_slides]
        
        if invalid_slides:
            error_msg = f"Invalid slide numbers: {invalid_slides} (valid range: 1-{max_slides})"
            logger.error(error_msg)
            result["error"] = error_msg
            return result
        
        # Step 4: Determine which slides to delete
        slides_to_keep_0based = set(s - 1 for s in slides_to_keep)  # Convert to 0-based
        slides_to_delete = [i for i in range(result["original_slide_count"]) 
                           if i not in slides_to_keep_0based]
        
        # Sort in reverse order to delete from end (prevents index shifting issues)
        slides_to_delete.sort(reverse=True)
        result["slides_deleted"] = len(slides_to_delete)
        
        logger.info(f"Will delete {len(slides_to_delete)} slides: {[i+1 for i in slides_to_delete]}")
        
        # Step 5: Delete unwanted slides using proven method
        deletion_errors = 0
        
        for slide_index in slides_to_delete:
            try:
                logger.debug(f"Deleting slide at index {slide_index} (slide number {slide_index + 1})")
                
                # Use python-pptx's internal structure to safely remove slides
                rId = presentation.slides._sldIdLst[slide_index].rId
                presentation.part.drop_rel(rId)
                del presentation.slides._sldIdLst[slide_index]
                
                logger.debug(f"Successfully deleted slide {slide_index + 1}")
                
            except Exception as delete_error:
                deletion_errors += 1
                logger.error(f"Error deleting slide {slide_index + 1}: {delete_error}")
                continue
        
        # Step 6: Save the modified presentation
        logger.debug("Saving modified presentation")
        presentation.save(str(temp_file))
        result["final_slide_count"] = len(presentation.slides)
        
        # Step 7: Verify the result
        logger.debug("Verifying saved presentation")
        try:
            verify_presentation = Presentation(str(temp_file))
            verify_count = len(verify_presentation.slides)
            
            if verify_count == result["final_slide_count"]:
                logger.debug(f"✅ Verification successful: {verify_count} slides")
            else:
                logger.warning(f"⚠️ Slide count mismatch: expected {result['final_slide_count']}, got {verify_count}")
                
        except Exception as verify_error:
            logger.error(f"❌ Verification failed: {verify_error}")
            result["error"] = f"Template verification failed: {verify_error}"
            return result
        
        # Step 8: Finalize successful result
        result["success"] = True
        result["template_path"] = str(temp_file)
        result["processing_time"] = time.time() - start_time
        
        logger.info(f"✅ Template creation successful!")
        logger.info(f"   Original slides: {result['original_slide_count']}")
        logger.info(f"   Final slides: {result['final_slide_count']}")
        logger.info(f"   Deleted slides: {result['slides_deleted']}")
        logger.info(f"   Processing time: {result['processing_time']:.2f} seconds")
        logger.info(f"   Template saved: {temp_file.name}")
        
        if deletion_errors > 0:
            logger.warning(f"   Deletion errors: {deletion_errors}")
        
    except Exception as critical_error:
        result["processing_time"] = time.time() - start_time
        result["error"] = str(critical_error)
        logger.error(f"❌ Critical error in template creation: {critical_error}")
        logger.debug(f"Template creation traceback: {traceback.format_exc()}")
    
    return result

def create_individual_template_files(template_results: List[Dict], output_dir: Path) -> Dict:
    """Create individual template files for download instead of merging."""
    logger = st.session_state.logger
    
    # Initialize result
    processing_result = {
        "success": False,
        "total_templates": len(template_results),
        "successful_templates": 0,
        "individual_files": [],
        "total_slides": 0,
        "processing_time": 0,
        "error": None,
        "timestamp": datetime.datetime.now().isoformat()
    }
    
    start_time = time.time()
    
    try:
        # Filter successful templates
        successful_templates = [r for r in template_results if r["success"] and r["template_path"]]
        processing_result["successful_templates"] = len(successful_templates)
        
        if not successful_templates:
            processing_result["error"] = "No successful templates to process"
            logger.error("No successful templates available")
            return processing_result
        
        logger.info(f"Creating {len(successful_templates)} individual template files...")
        
        # Generate timestamp for consistent naming
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Process each successful template
        for i, template_result in enumerate(successful_templates):
            try:
                source_path = template_result["template_path"]
                source_file_name = template_result["source_file"]
                
                # Generate clean filename
                base_name = os.path.splitext(source_file_name)[0]
                # Remove any problematic characters
                clean_name = re.sub(r'[^\w\-_\.]', '_', base_name)
                final_filename = f"{clean_name}_assembled_{timestamp}.pptx"
                final_file_path = output_dir / final_filename
                
                # Copy template to final location with proper name
                shutil.copy2(source_path, final_file_path)
                
                # Verify file and get slide count
                final_presentation = Presentation(str(final_file_path))
                slide_count = len(final_presentation.slides)
                file_size = final_file_path.stat().st_size
                
                # Add to results
                file_info = {
                    "original_name": source_file_name,
                    "final_filename": final_filename,
                    "file_path": str(final_file_path),
                    "slide_count": slide_count,
                    "file_size_mb": file_size / (1024 * 1024),
                    "slides_kept": template_result["slides_requested"].copy(),
                    "processing_successful": True
                }
                
                processing_result["individual_files"].append(file_info)
                processing_result["total_slides"] += slide_count
                
                logger.info(f"✅ Created: {final_filename} ({slide_count} slides, {file_size/1024/1024:.1f}MB)")
                
            except Exception as file_error:
                logger.error(f"❌ Error processing template {i+1}: {file_error}")
                
                # Add failed file info
                file_info = {
                    "original_name": template_result.get("source_file", f"Template_{i+1}"),
                    "final_filename": None,
                    "file_path": None,
                    "slide_count": 0,
                    "file_size_mb": 0,
                    "slides_kept": template_result.get("slides_requested", []),
                    "processing_successful": False,
                    "error": str(file_error)
                }
                
                processing_result["individual_files"].append(file_info)
                continue
        
        # Check if any files were successfully created
        successful_files = [f for f in processing_result["individual_files"] if f["processing_successful"]]
        
        if successful_files:
            processing_result["success"] = True
            processing_result["processing_time"] = time.time() - start_time
            
            logger.info(f"✅ Template processing completed successfully!")
            logger.info(f"   Files created: {len(successful_files)}")
            logger.info(f"   Total slides: {processing_result['total_slides']}")
            logger.info(f"   Processing time: {processing_result['processing_time']:.2f} seconds")
            
        else:
            processing_result["error"] = "No template files could be created successfully"
            logger.error("❌ All template processing failed")
            
    except Exception as critical_error:
        processing_result["processing_time"] = time.time() - start_time
        processing_result["error"] = str(critical_error)
        logger.error(f"❌ Critical error in template processing: {critical_error}")
        logger.debug(f"Template processing traceback: {traceback.format_exc()}")
    
    return processing_result

# ============================================================================
# STREAMLIT USER INTERFACE
# ============================================================================

def initialize_streamlit_session():
    """Initialize Streamlit session state variables and logging."""
    
    # Generate unique session ID for this run
    if 'session_id' not in st.session_state:
        st.session_state.session_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
    
    # Initialize logging
    if 'logger' not in st.session_state:
        log_file, logger = setup_comprehensive_logging(st.session_state.session_id)
        st.session_state.log_file = log_file
        st.session_state.logger = logger
    
    # Initialize other session variables
    if 'uploaded_files_info' not in st.session_state:
        st.session_state.uploaded_files_info = []
    
    if 'llm_client' not in st.session_state:
        st.session_state.llm_client = None
    
    if 'analysis_results' not in st.session_state:
        st.session_state.analysis_results = []
    
    if 'processing_results' not in st.session_state:
        st.session_state.processing_results = {}

def render_header():
    """Render the application header and configuration."""
    
    st.set_page_config(
        page_title=f"{APP_NAME} v{APP_VERSION}",
        page_icon="📊",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    st.title(f"📊 {APP_NAME}")
    st.caption(f"Version {APP_VERSION} - Template Duplication + AI Analysis")
    
    # Header information
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.info(f"**Session ID:** `{st.session_state.session_id}`")
    
    with col2:
        st.info(f"**LLM Model:** `{DEEPSEEK_MODEL}`")
    
    with col3:
        st.info(f"**Debug Mode:** `{'ON' if DEBUG_MODE else 'OFF'}`")
    
    st.divider()

def render_api_configuration():
    """Render API configuration section - FIXED: NO NESTED EXPANDERS."""
    
    st.subheader("🔑 API Configuration")
    
    # THE ONLY EXPANDER IN THIS ENTIRE FUNCTION
    with st.expander("OpenRouter API Setup", expanded=not st.session_state.llm_client):
        st.markdown("""
        **DeepSeek R1 via OpenRouter Configuration**
        
        This application uses DeepSeek R1 through OpenRouter for intelligent slide analysis.
        The API key is securely configured and ready to use.
        
        **External AI Service Notice:** This application processes your content through external AI services.
        Please only upload non-confidential content.
        """)
        
        # API Key handling - now uses Streamlit secrets
        try:
            api_key_input = st.secrets["OPENROUTER_API_KEY"]
            st.success("🔑 **API Key loaded from secure configuration**")
            st.info("✅ Ready for AI analysis - no additional setup required!")
        except KeyError:
            st.error("❌ **API Key not configured**")
            st.error("Please contact the administrator to configure the OpenRouter API key.")
            api_key_input = None
        except Exception as e:
            st.error(f"❌ **Error accessing API key**: {str(e)}")
            api_key_input = None
        
        # SSL Configuration for Corporate Networks
        st.markdown("**🔒 Corporate Network Settings**")
        bypass_ssl = st.checkbox(
            "Bypass SSL Certificate Verification",
            value=False,
            help="Enable this if you're on a corporate network with SSL inspection. This is safe for API calls but should only be used if needed."
        )
        
        if bypass_ssl:
            st.warning("⚠️ **SSL verification disabled** - Only use this for corporate networks that intercept HTTPS traffic")
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            connect_button = st.button("🔗 Connect to DeepSeek API", type="primary")
        
        with col2:
            test_button = st.button("🧪 Test Basic API", help="Test OpenRouter API access")
        
        # Basic API test (without model call)
        if test_button and api_key_input:
            with st.spinner("Testing basic OpenRouter API access..."):
                try:
                    try:
                        import requests
                    except ImportError:
                        st.error("Missing requests dependency. Please install with: pip install requests")
                        return
                    
                    # Configure requests for SSL
                    session = requests.Session()
                    if bypass_ssl:
                        session.verify = False
                        try:
                            import urllib3
                            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
                        except ImportError:
                            pass
                        st.warning("🔒 SSL verification bypassed for this test")
                    
                    # Simple test request to OpenRouter
                    test_url = f"{OPENROUTER_BASE_URL}/models"
                    headers = {
                        "Authorization": f"Bearer {api_key_input}",
                        "HTTP-Referer": "https://slide-assembler.internal",
                        "X-Title": "PowerPoint Slide Assembler Test"
                    }
                    
                    st.session_state.logger.info(f"Testing basic API connectivity (SSL bypass: {bypass_ssl})...")
                    response = session.get(test_url, headers=headers, timeout=10)
                    
                    if response.status_code == 200:
                        st.success("✅ Basic OpenRouter API connectivity successful!")
                        models_data = response.json()
                        
                        # Check if DeepSeek model is available
                        available_models = [model.get('id', '') for model in models_data.get('data', [])]
                        deepseek_available = any(DEEPSEEK_MODEL in model_id for model_id in available_models)
                        
                        if deepseek_available:
                            st.success(f"✅ DeepSeek model ({DEEPSEEK_MODEL}) is available!")
                        else:
                            st.warning(f"⚠️ DeepSeek model ({DEEPSEEK_MODEL}) not found in available models")
                            
                            # Show some available models
                            deepseek_models = [m for m in available_models if 'deepseek' in m.lower()]
                            if deepseek_models:
                                st.info(f"Available DeepSeek models: {', '.join(deepseek_models[:3])}")
                            
                            st.info(f"Total available models: {len(available_models)}")
                        
                    elif response.status_code == 401:
                        st.error("❌ API key is invalid or unauthorized")
                    elif response.status_code == 403:
                        st.error("❌ Access forbidden - check your API key permissions")
                    else:
                        st.error(f"❌ API test failed: HTTP {response.status_code}")
                        st.code(response.text[:500])
                        
                except requests.exceptions.SSLError as ssl_err:
                    st.error("🔒 **SSL Certificate Error** - Try enabling 'Bypass SSL Certificate Verification' above")
                    st.code(str(ssl_err)[:300])
                except requests.exceptions.ConnectTimeout:
                    st.error("🌐 Connection timeout - check your internet connection or firewall settings")
                except requests.exceptions.ConnectionError as conn_err:
                    if "certificate verify failed" in str(conn_err).lower():
                        st.error("🔒 **SSL Certificate Error** - Try enabling 'Bypass SSL Certificate Verification' above")
                    else:
                        st.error("🌐 Connection error - OpenRouter may be unreachable from your network")
                    st.code(str(conn_err)[:300])
                except Exception as e:
                    st.error(f"❌ Test failed: {str(e)}")
                    st.code(str(e)[:300])
        
        # Full connection attempt with DeepSeek model
        if connect_button:
            if api_key_input:
                with st.spinner("Testing full API connection with DeepSeek model..."):
                    client = initialize_llm_client(api_key_input, bypass_ssl)
                    if client:
                        st.session_state.llm_client = client
                        st.session_state.api_key = api_key_input
                        st.session_state.bypass_ssl = bypass_ssl
                        st.success("✅ Successfully connected to DeepSeek API!")
                        st.balloons()
                        st.rerun()
                    else:
                        st.error("❌ Failed to connect to API. Check the error details above.")
                        
                        # Suggest trying SSL bypass if not already enabled
                        if not bypass_ssl and "certificate" in st.session_state.get('last_error', '').lower():
                            st.info("💡 **Suggestion:** Try enabling 'Bypass SSL Certificate Verification' above and reconnecting")
            else:
                st.warning("API key not available. Please contact the administrator.")
        
        # Current status
        if st.session_state.llm_client:
            ssl_status = "🔒 Bypassed" if st.session_state.get('bypass_ssl', False) else "🔒 Verified"
            st.success(f"🟢 **DeepSeek API Connected** - AI analysis enabled! (SSL: {ssl_status})")
            
            # Option to disconnect
            if st.button("🔌 Disconnect API"):
                st.session_state.llm_client = None
                st.session_state.api_key = None
                st.session_state.bypass_ssl = False
                st.rerun()
    
    # END OF FUNCTION - ABSOLUTELY NO MORE EXPANDERS!

def render_troubleshooting_help():
    """Render troubleshooting help as regular content (NO EXPANDERS)."""
    
    if not st.session_state.llm_client:
        st.markdown("---")
        st.markdown("### 🏢 Corporate Network Troubleshooting")
        
        st.markdown("""
        **SSL Certificate Issues? Here's how to fix it:**
        
        🔒 **SSL Certificate Problem:**
        Your corporate firewall is intercepting HTTPS traffic and presenting its own certificate.
        
        **Quick Fix:**
        1. ✅ **Enable "Bypass SSL Certificate Verification" above**
        2. 🔗 **Click "Connect to DeepSeek API"**
        3. 🎉 **Should work immediately**
        
        🏢 **Permanent Solutions:**
        - **Ask IT to whitelist:** `openrouter.ai` and `*.openrouter.ai`
        - **Request direct access** to AI APIs for business tools
        - **Install corporate certificates** in Python environment
        
        📱 **Alternative Testing:**
        - **Personal hotspot:** Test outside corporate network
        - **VPN:** Use personal VPN if allowed
        - **Different network:** Try from home/café
        
        ⚠️ **Security Note:** Bypassing SSL is safe for this API usage but should only be used when necessary.
        """)
        
        # Quick tests in columns
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("🔍 Test SSL Configuration"):
                with st.spinner("Testing SSL connectivity to OpenRouter..."):
                    try:
                        import ssl
                        import socket
                        
                        # Test SSL connection to OpenRouter
                        context = ssl.create_default_context()
                        
                        with socket.create_connection(("openrouter.ai", 443), timeout=10) as sock:
                            with context.wrap_socket(sock, server_hostname="openrouter.ai") as ssock:
                                st.success("✅ SSL connection to OpenRouter successful!")
                                st.info(f"SSL Version: {ssock.version()}")
                                
                    except ssl.SSLError as ssl_err:
                        st.error("🔒 SSL Error - Corporate certificate interception detected")
                        st.code(str(ssl_err))
                        st.info("💡 Solution: Enable 'Bypass SSL Certificate Verification' above")
                    except Exception as e:
                        st.error(f"❌ SSL test failed: {str(e)}")
        
        with col2:
            if st.button("🌐 Test Basic Network"):
                with st.spinner("Testing network connectivity..."):
                    try:
                        try:
                            import requests
                        except ImportError:
                            st.error("Missing requests dependency. Please install with: pip install requests")
                            return
                        
                        # Test basic internet
                        response = requests.get("https://httpbin.org/ip", timeout=5)
                        if response.status_code == 200:
                            st.success("✅ Basic internet connectivity working")
                            
                            # Test OpenRouter domain
                            try:
                                response = requests.get("https://openrouter.ai", timeout=10)
                                if response.status_code == 200:
                                    st.success("✅ OpenRouter domain is reachable")
                                else:
                                    st.warning(f"⚠️ OpenRouter domain returned status {response.status_code}")
                            except requests.exceptions.SSLError:
                                st.warning("🔒 SSL issues detected with OpenRouter domain")
                                st.info("💡 Try enabling SSL bypass above")
                        else:
                            st.error("❌ Basic internet connectivity failed")
                            
                    except Exception as e:
                        st.error(f"❌ Network test failed: {str(e)}")

def render_debug_info():
    """Render debug information as regular content (NO EXPANDERS)."""
    
    if DEBUG_MODE:
        st.markdown("---")
        st.markdown("### 🐛 Debug Information")
        
        debug_info = f"""
API Configuration Debug Info:
- Base URL: {OPENROUTER_BASE_URL}
- Model: {DEEPSEEK_MODEL}
- API Key Source: Streamlit Secrets
- SSL Bypass: {st.session_state.get('bypass_ssl', False)}
- Client Status: {'Connected' if st.session_state.llm_client else 'Not Connected'}
        """
        
        st.code(debug_info)

def render_file_upload():
    """Render file upload section with privacy protection."""
    
    st.subheader("📁 File Upload")
    
    # File upload widget
    uploaded_files = st.file_uploader(
        "Upload PowerPoint Files (.pptx)",
        type=['pptx'],
        accept_multiple_files=True,
        help=f"Upload up to {MAX_TOTAL_FILES} PowerPoint files, max {MAX_FILE_SIZE_MB}MB each"
    )
    
    # Show privacy warning if files are uploaded
    if uploaded_files:
        
        # Show privacy warning and block functionality until acknowledged
        if not show_privacy_warning():
            st.info("👆 **Please acknowledge the privacy warning above to proceed with file processing.**")
            return
        
        # Privacy acknowledged - proceed with file processing
        st.session_state.uploaded_files_info = []
        
        # Add privacy reminder
        st.success("🔒 **Privacy acknowledged** - Processing uploaded files...")
        
        # Process uploaded files
        for uploaded_file in uploaded_files:
            file_size_mb = uploaded_file.size / (1024 * 1024)
            
            if file_size_mb > MAX_FILE_SIZE_MB:
                st.error(f"❌ {uploaded_file.name} is too large ({file_size_mb:.1f}MB). Max size: {MAX_FILE_SIZE_MB}MB")
                continue
            
            # Save file temporarily and analyze
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(uploaded_file.getbuffer())
                tmp_file_path = tmp_file.name
            
            try:
                # Get slide count
                pres = Presentation(tmp_file_path)
                slide_count = len(pres.slides)
                
                file_info = {
                    "name": uploaded_file.name,
                    "size_mb": file_size_mb,
                    "slide_count": slide_count,
                    "temp_path": tmp_file_path,
                    "upload_time": datetime.datetime.now().isoformat()
                }
                
                st.session_state.uploaded_files_info.append(file_info)
                st.session_state.logger.info(f"File uploaded: {uploaded_file.name} ({slide_count} slides, {file_size_mb:.1f}MB)")
                
            except Exception as e:
                st.error(f"❌ Error reading {uploaded_file.name}: {str(e)}")
                os.unlink(tmp_file_path)
        
        # Display uploaded files summary
        if st.session_state.uploaded_files_info:
            st.success(f"✅ {len(st.session_state.uploaded_files_info)} files uploaded successfully!")
            
            # Privacy reminder
            st.info("🔒 **Remember**: Only non-confidential content should be uploaded for AI analysis")
            
            # Files summary table
            files_data = []
            total_slides = 0
            total_size = 0
            
            for file_info in st.session_state.uploaded_files_info:
                files_data.append({
                    "File Name": file_info["name"],
                    "Slides": file_info["slide_count"],
                    "Size (MB)": f"{file_info['size_mb']:.1f}"
                })
                total_slides += file_info["slide_count"]
                total_size += file_info["size_mb"]
            
            st.dataframe(files_data, use_container_width=True)
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Total Files", len(st.session_state.uploaded_files_info))
            with col2:
                st.metric("Total Slides", total_slides)
            with col3:
                st.metric("Total Size (MB)", f"{total_size:.1f}")
    
    # Show privacy status in session
    elif 'privacy_acknowledged' in st.session_state and st.session_state.privacy_acknowledged:
        st.info("🔒 **Privacy settings acknowledged** - Ready for file upload")
        
        # Option to reset privacy acknowledgment
        if st.button("🔄 Reset Privacy Acknowledgment", help="Click if you want to see the privacy warning again"):
            st.session_state.privacy_acknowledged = False
            st.rerun()

def render_content_analysis():
    """Render content analysis section with LLM integration."""
    
    if not st.session_state.uploaded_files_info:
        st.info("👆 Please upload PowerPoint files first to enable content analysis.")
        return
    
    st.subheader("🧠 AI-Powered Content Analysis")
    
    if not st.session_state.llm_client:
        st.warning("⚠️ Connect to DeepSeek API first to enable intelligent content analysis.")
        return
    
    st.markdown("**Content Analysis Configuration**")
    
    # Key message input
    key_message = st.text_area(
        "🎯 Key Message for Your Presentation",
        placeholder="Enter the main message or objective for your presentation...",
        help="This helps the AI understand your goal and provide relevant recommendations.",
        height=100
    )
    
    # Analysis type selection
    col1, col2 = st.columns([2, 1])
    
    with col1:
        analysis_type = st.selectbox(
            "Analysis Type",
            ["comprehensive", "relevance", "quick"],
            help="Choose the depth of analysis you want"
        )
    
    with col2:
        # Content optimization toggle
        optimize_content = st.checkbox(
            "Optimize for Large Content",
            help="Automatically truncate very large content to avoid timeouts",
            value=True
        )
    
    # Show analysis type descriptions
    analysis_descriptions = {
        "comprehensive": "🔍 **Comprehensive**: Full detailed analysis (2-5 minutes for large content)",
        "relevance": "🎯 **Relevance**: Focus on content relevance to key message (1-2 minutes)",
        "quick": "⚡ **Quick**: Fast overview and top recommendations (30-60 seconds)"
    }
    
    st.info(analysis_descriptions[analysis_type])
    
    # File selection for analysis
    file_options = {info["name"]: info for info in st.session_state.uploaded_files_info}
    selected_files = st.multiselect(
        "Select Files to Analyze",
        options=list(file_options.keys()),
        default=list(file_options.keys()),
        help="Choose which files to include in the content analysis"
    )
    
    if st.button("🔍 Analyze Content with AI", type="primary", disabled=not key_message or not selected_files):
        
        # Pre-analysis content size check
        total_content_size = 0
        for file_name in selected_files:
            file_info = file_options[file_name]
            # Quick content size estimate (rough)
            estimated_size = file_info["slide_count"] * 450  # ~450 chars per slide average
            total_content_size += estimated_size
        
        # Show size warning if needed
        if total_content_size > 50000:
            st.warning(f"⚠️ **Large Content Warning**: Estimated {total_content_size:,} characters")
            st.info("💡 **For faster processing**, consider:")
            st.info("• Using 'quick' or 'relevance' analysis instead of 'comprehensive'")
            st.info("• Selecting fewer files or slides")
            
        elif total_content_size > 20000:
            st.info(f"📊 **Large Content**: Estimated {total_content_size:,} characters - may take 2-3 minutes")
        
        # Progress placeholder
        progress_placeholder = st.empty()
        
        with st.spinner("🤖 AI is analyzing your presentation content..."):
            
            # Step 1: Extract content
            progress_placeholder.info("📄 Step 1/2: Extracting text content from slides...")
            
            # Extract content from selected files
            all_content = ""
            extraction_stats = {"total_files": 0, "total_slides": 0, "total_characters": 0}
            
            for i, file_name in enumerate(selected_files):
                file_info = file_options[file_name]
                
                progress_placeholder.info(f"📄 Extracting from file {i+1}/{len(selected_files)}: {file_name}")
                st.session_state.logger.info(f"Extracting content from: {file_name}")
                
                # Extract text content
                extraction_result = extract_slide_text_content(file_info["temp_path"])
                
                if extraction_result["success"]:
                    all_content += f"\\n\\n=== FILE: {file_name} ===\\n"
                    all_content += extraction_result["combined_content"]
                    
                    stats = extraction_result["extraction_stats"]
                    extraction_stats["total_files"] += 1
                    extraction_stats["total_slides"] += stats["processed_slides"]
                    extraction_stats["total_characters"] += stats["total_characters"]
                    
                    st.success(f"✅ Content extracted from {file_name}: {stats['processed_slides']} slides")
                else:
                    st.error(f"❌ Failed to extract content from {file_name}: {extraction_result.get('error', 'Unknown error')}")
            
            # Step 2: Perform LLM analysis
            if all_content.strip():
                actual_size = extraction_stats['total_characters']
                
                # Content optimization for large files
                optimized_content = all_content
                optimization_applied = False
                
                if optimize_content and actual_size > 40000:
                    st.warning(f"⚠️ **Large content detected**: {actual_size:,} characters")
                    st.info("🔧 **Applying content optimization** to prevent timeouts...")
                    
                    # Truncate content to reasonable size while preserving structure
                    max_chars = 35000 if analysis_type == "comprehensive" else 25000
                    
                    if actual_size > max_chars:
                        # Keep file headers and truncate content proportionally
                        lines = optimized_content.split('\n')
                        file_headers = [line for line in lines if line.startswith('=== FILE:')]
                        
                        # Calculate how much content per file
                        num_files = len(file_headers)
                        chars_per_file = max_chars // num_files if num_files > 0 else max_chars
                        
                        optimized_lines = []
                        current_chars = 0
                        
                        for line in lines:
                            if current_chars + len(line) > max_chars:
                                optimized_lines.append(f"\n[CONTENT TRUNCATED FOR OPTIMIZATION - ORIGINAL SIZE: {actual_size:,} CHARS]")
                                break
                            optimized_lines.append(line)
                            current_chars += len(line)
                        
                        optimized_content = '\n'.join(optimized_lines)
                        optimization_applied = True
                        
                        st.info(f"✂️ **Content optimized**: {actual_size:,} → {len(optimized_content):,} characters")
                        st.session_state.logger.info(f"Content optimized: {actual_size:,} → {len(optimized_content):,} characters")
                
                progress_placeholder.info(f"🧠 Step 2/2: AI analyzing {len(optimized_content):,} characters...")
                
                st.info(f"📊 Processing {len(optimized_content):,} characters from {extraction_stats['total_slides']} slides across {extraction_stats['total_files']} files")
                
                if optimization_applied:
                    st.info("⚡ **Optimization applied** - Analysis should complete faster")
                
                # Show estimated processing time
                content_size = len(optimized_content)
                if content_size > 50000:
                    st.warning("⏳ **Large content** - Analysis may take 3-5 minutes")
                elif content_size > 20000:
                    st.info("⏳ **Processing time**: Estimated 1-3 minutes")
                else:
                    st.info("⏳ **Processing time**: Estimated 30-60 seconds")
                
                analysis_result = analyze_slide_content_with_llm(
                    st.session_state.llm_client,
                    optimized_content,  # Use optimized content
                    key_message,
                    analysis_type
                )
                
                # Clear progress
                progress_placeholder.empty()
                
                if analysis_result["success"]:
                    st.success(f"✅ AI analysis completed in {analysis_result['response_time']:.2f} seconds!")
                    
                    # Store results with optimization info
                    analysis_result["key_message"] = key_message
                    analysis_result["selected_files"] = selected_files
                    analysis_result["extraction_stats"] = extraction_stats
                    analysis_result["optimization_applied"] = optimization_applied
                    if optimization_applied:
                        analysis_result["original_size"] = actual_size
                        analysis_result["optimized_size"] = len(optimized_content)
                    
                    st.session_state.analysis_results.append(analysis_result)
                    
                    # Display results
                    st.subheader("🎯 AI Analysis Results")
                    
                    # Show optimization notice if applied
                    if optimization_applied:
                        st.info(f"ℹ️ **Content was optimized** from {actual_size:,} to {len(optimized_content):,} characters for faster processing")
                    
                    st.markdown(analysis_result["analysis_text"])
                    
                    # Show performance stats
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Processing Time", f"{analysis_result['response_time']:.1f}s")
                    with col2:
                        if analysis_result.get("token_usage") and analysis_result["token_usage"]["total_tokens"]:
                            st.metric("Tokens Used", f"{analysis_result['token_usage']['total_tokens']:,}")
                        else:
                            display_size = len(optimized_content) if optimization_applied else actual_size
                            st.metric("Content Size", f"{display_size:,} chars")
                    with col3:
                        optimization_status = "Optimized" if optimization_applied else "Original"
                        st.metric("Content", optimization_status)
                    
                    # Show token usage details if available
                    if analysis_result.get("token_usage") and analysis_result["token_usage"]["total_tokens"]:
                        usage = analysis_result["token_usage"]
                        st.caption(f"💰 Token breakdown: {usage['prompt_tokens']:,} prompt + {usage['completion_tokens']:,} completion = {usage['total_tokens']:,} total")
                    
                    # Show content stats
                    if optimization_applied:
                        st.caption(f"📊 Content: {actual_size:,} chars extracted → {len(optimized_content):,} chars analyzed ({analysis_type} analysis)")
                    else:
                        st.caption(f"📊 Content: {actual_size:,} characters analyzed ({analysis_type} analysis)")
                    
                else:
                    st.error(f"❌ AI analysis failed: {analysis_result.get('error', 'Unknown error')}")
                    
                    # Show detailed error information
                    error_type = analysis_result.get('error_type', 'Unknown')
                    content_stats = analysis_result.get('content_stats', {})
                    
                    st.error(f"**Error Type**: {error_type}")
                    
                    # Show troubleshooting based on error type
                    if 'timeout' in analysis_result.get('error', '').lower():
                        st.warning("🕐 **Timeout Error** - The request took too long to complete")
                        st.info("**💡 Try these solutions:**")
                        st.info("• Switch to 'quick' or 'relevance' analysis")
                        st.info("• Select fewer slides to reduce content size")
                        st.info("• Try again in a few minutes")
                        st.info(f"• Current content: {content_stats.get('input_characters', 0):,} characters")
                        
                    elif 'rate' in analysis_result.get('error', '').lower():
                        st.warning("🚦 **Rate Limit** - Too many requests")
                        st.info("• Wait 1-2 minutes before trying again")
                        st.info("• Consider upgrading your OpenRouter plan")
                        
                    else:
                        st.info("**Troubleshooting suggestions:**")
                        troubleshooting = analysis_result.get('troubleshooting', {})
                        for action in troubleshooting.get('suggested_actions', []):
                            st.info(f"• {action}")
                    
                    # Debug information
                    if DEBUG_MODE:
                        st.markdown("**🐛 Debug Information:**")
                        st.code(f"""
Error Details:
- Error Type: {error_type}
- Content Size: {content_stats.get('input_characters', 0):,} characters
- Estimated Tokens: {content_stats.get('estimated_input_tokens', 0):,}
- Analysis Type: {analysis_type}
- Model: {DEEPSEEK_MODEL}
- Files: {len(selected_files)}
- Slides: {extraction_stats['total_slides']}
                        """)
            else:
                progress_placeholder.empty()
                st.warning("⚠️ No content extracted from selected files.")

def render_slide_assembly():
    """Render slide assembly section with template duplication."""
    
    if not st.session_state.uploaded_files_info:
        st.info("👆 Please upload PowerPoint files first to enable slide assembly.")
        return
    
    st.subheader("🔧 Slide Assembly (Template Duplication)")
    
    st.markdown("""
    **Template Duplication Method** - Preserves ALL formatting, themes, embedded objects, and media.
    This proven approach creates perfect copies of your original slides.
    """)
    
    # Multi-deck limitation disclaimer
    if len(st.session_state.uploaded_files_info) > 1:
        st.warning("""
        **📋 Multi-Deck Notice**: When multiple PowerPoint files are selected, this tool will create **individual assembled files** for each deck rather than combining them into one file.
        
        **Why?** Merging multiple PowerPoint files while preserving all formatting, themes, and layouts is complex and can break formatting.
        
        **Result**: You'll get separate `.pptx` files that you can manually combine in PowerPoint if needed.
        """)
    
    st.markdown("**Assembly Configuration**")
    
    # File and slide selection
    assembly_configs = []
    
    for i, file_info in enumerate(st.session_state.uploaded_files_info):
        st.markdown(f"### 📄 {file_info['name']}")
        st.caption(f"Available slides: 1-{file_info['slide_count']} ({file_info['slide_count']} total)")
        
        # Slide range input
        slide_range = st.text_input(
            f"Slide ranges for {file_info['name']}",
            placeholder="e.g., 1-3,6,9-12 or 1,2,3,5",
            help="Enter slide numbers or ranges to include in final presentation",
            key=f"slide_range_{i}"
        )
        
        if slide_range.strip():
            try:
                parsed_slides = parse_slide_ranges(slide_range)
                
                # Validate slide numbers
                max_slides = file_info['slide_count']
                invalid_slides = [s for s in parsed_slides if s < 1 or s > max_slides]
                
                if invalid_slides:
                    st.error(f"❌ Invalid slide numbers: {invalid_slides} (valid range: 1-{max_slides})")
                else:
                    st.success(f"✅ Valid slides selected: {parsed_slides} ({len(parsed_slides)} slides)")
                    
                    assembly_configs.append({
                        "file_info": file_info,
                        "slide_range": slide_range,
                        "parsed_slides": parsed_slides
                    })
                    
            except Exception as e:
                st.error(f"❌ Error parsing slide range: {str(e)}")
        
        st.divider()
    
    # Assembly button
    if assembly_configs:
        total_slides = sum(len(config["parsed_slides"]) for config in assembly_configs)
        num_files = len(assembly_configs)
        
        if num_files == 1:
            st.info(f"📊 Ready to assemble: {total_slides} slides from 1 file → **1 output file**")
        else:
            st.info(f"📊 Ready to assemble: {total_slides} slides from {num_files} files → **{num_files} individual output files**")
        
        if st.button("🚀 Create Assembled Presentation(s)", type="primary"):
            
            with st.spinner("🔄 Creating assembled presentations using template duplication..."):
                
                # Create output directory
                output_dir = Path(OUTPUT_DIR_NAME)
                output_dir.mkdir(exist_ok=True)
                
                # Process each file configuration
                template_results = []
                
                for i, config in enumerate(assembly_configs):
                    file_info = config["file_info"]
                    slides_to_keep = config["parsed_slides"]
                    
                    st.session_state.logger.info(f"Processing file {i+1}/{len(assembly_configs)}: {file_info['name']}")
                    
                    # Create trimmed template
                    template_result = create_trimmed_template(
                        file_info["temp_path"],
                        slides_to_keep,
                        output_dir,
                        f"file_{i+1}"
                    )
                    
                    template_results.append(template_result)
                    
                    if template_result["success"]:
                        st.success(f"✅ Template created for {file_info['name']}: {template_result['final_slide_count']} slides")
                    else:
                        st.error(f"❌ Failed to create template for {file_info['name']}: {template_result.get('error', 'Unknown error')}")
                
                # Create individual template files
                st.info("📁 Creating individual presentation files...")
                
                processing_result = create_individual_template_files(template_results, output_dir)
                
                if processing_result["success"]:
                    st.success(f"🎉 **Assembly completed successfully!**")
                    
                    # Display overall results
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        st.metric("Files Created", len([f for f in processing_result["individual_files"] if f["processing_successful"]]))
                    with col2:
                        st.metric("Total Slides", processing_result["total_slides"])
                    with col3:
                        st.metric("Processing Time", f"{processing_result['processing_time']:.1f}s")
                    
                    # Display individual file results and download buttons
                    st.subheader("📥 Download Your Assembled Presentations")
                    
                    successful_files = [f for f in processing_result["individual_files"] if f["processing_successful"]]
                    
                    if successful_files:
                        for i, file_info in enumerate(successful_files):
                            with st.container():
                                st.markdown(f"### 📄 {file_info['original_name']}")
                                
                                col1, col2, col3 = st.columns([2, 1, 1])
                                
                                with col1:
                                    st.markdown(f"**Output:** `{file_info['final_filename']}`")
                                    st.caption(f"Slides included: {', '.join(map(str, file_info['slides_kept']))}")
                                
                                with col2:
                                    st.metric("Slides", file_info['slide_count'])
                                    st.metric("Size", f"{file_info['file_size_mb']:.1f} MB")
                                
                                with col3:
                                    # Download button for this file
                                    file_path = Path(file_info['file_path'])
                                    if file_path.exists():
                                        with open(file_path, "rb") as file:
                                            st.download_button(
                                                label=f"📥 Download",
                                                data=file.read(),
                                                file_name=file_info['final_filename'],
                                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                                type="primary",
                                                key=f"download_{i}",
                                                use_container_width=True
                                            )
                                        
                                        st.caption(f"📁 Saved: `{file_path.absolute()}`")
                                
                                if i < len(successful_files) - 1:
                                    st.divider()
                        
                        # Summary message
                        if len(successful_files) > 1:
                            st.info(f"""
                            **📋 Multiple Files Created**: You now have {len(successful_files)} individual presentation files.
                            
                            **To combine them manually:**
                            1. Open the first presentation in PowerPoint
                            2. Use "Insert → Slides from Other Presentation" to add slides from other files
                            3. Arrange slides as needed
                            
                            **Why separate files?** This preserves all formatting, themes, and layouts perfectly.
                            """)
                    
                    # Show failed files if any
                    failed_files = [f for f in processing_result["individual_files"] if not f["processing_successful"]]
                    if failed_files:
                        st.error("❌ **Some files failed to process:**")
                        for file_info in failed_files:
                            st.error(f"• {file_info['original_name']}: {file_info.get('error', 'Unknown error')}")
                    
                    # Store results
                    st.session_state.processing_results = {
                        "template_results": template_results,
                        "processing_result": processing_result,
                        "timestamp": datetime.datetime.now().isoformat()
                    }
                    
                else:
                    st.error(f"❌ Assembly failed: {processing_result.get('error', 'Unknown error')}")
    else:
        st.info("👆 Configure slide ranges for your files to enable assembly.")

def render_session_info():
    """Render session information and logs in sidebar."""
    
    with st.sidebar:
        st.header("📋 Session Information")
        
        # Session details
        st.subheader("Session Details")
        st.text(f"ID: {st.session_state.session_id}")
        st.text(f"Started: {datetime.datetime.now().strftime('%H:%M:%S')}")
        st.text(f"Log File: {st.session_state.log_file.name}")
        
        # API status
        st.subheader("API Status")
        if st.session_state.llm_client:
            st.success("🟢 DeepSeek Connected")
        else:
            st.error("🔴 No LLM Connection")
        
        # File summary
        if st.session_state.uploaded_files_info:
            st.subheader("Uploaded Files")
            for file_info in st.session_state.uploaded_files_info:
                st.text(f"📄 {file_info['name']}")
                st.text(f"   {file_info['slide_count']} slides")
        
        # Analysis results summary
        if st.session_state.analysis_results:
            st.subheader("Analysis Results")
            st.text(f"Completed: {len(st.session_state.analysis_results)}")
        
        # Processing results summary
        if st.session_state.processing_results:
            st.subheader("Assembly Results")
            processing_result = st.session_state.processing_results.get("processing_result", {})
            if processing_result.get("success"):
                successful_files = len([f for f in processing_result.get("individual_files", []) if f.get("processing_successful", False)])
                total_slides = processing_result.get("total_slides", 0)
                st.success(f"✅ {successful_files} files, {total_slides} slides")
            else:
                st.error("❌ Assembly failed")
        
        # Debug information
        if DEBUG_MODE:
            st.subheader("🐛 Debug Info")
            st.text(f"Session state keys: {len(st.session_state.keys())}")
            st.text(f"Temp files: {len(st.session_state.uploaded_files_info)}")

# ============================================================================
# MAIN APPLICATION
# ============================================================================

def main():
    """Main application entry point."""
    
    # Initialize session
    initialize_streamlit_session()
    
    # Render UI components
    render_header()
    render_api_configuration()
    
    # Render troubleshooting and debug sections as regular content (NO EXPANDERS)
    render_troubleshooting_help()
    render_debug_info()
    
    st.divider()
    
    # Main content tabs
    tab1, tab2, tab3 = st.tabs(["📁 File Upload", "🧠 AI Analysis", "🔧 Assembly"])
    
    with tab1:
        render_file_upload()
    
    with tab2:
        render_content_analysis()
    
    with tab3:
        render_slide_assembly()
    
    # Sidebar information
    render_session_info()
    
    # Footer
    st.divider()
    st.caption(f"{APP_NAME} v{APP_VERSION} - Built with Streamlit & DeepSeek")

if __name__ == "__main__":
    main()
